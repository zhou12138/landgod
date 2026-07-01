#!/usr/bin/env python3
"""LandGod business report demo MCP.

A small business-facing demo for the Enterprise Execution Harness story:
- mock ERP order data
- mock Finance invoice data
- credential-aware trusted tools
- CSV/HTML/PPTX report generation

The server intentionally returns credential metadata only; it never returns secret values.
"""
from __future__ import annotations

import argparse
import csv
import json
import os
from collections import defaultdict
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import anyio
from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp.types import Tool

ROOT = Path(__file__).resolve().parent
DATA_DIR = ROOT / "data"
DEFAULT_OUTPUT_DIR = ROOT / "output"

server = Server("business-report-demo")


def _read_csv(path: Path) -> list[dict[str, str]]:
    with path.open("r", encoding="utf-8", newline="") as f:
        return list(csv.DictReader(f))


def _read_targets(month: str) -> dict[str, Any]:
    path = DATA_DIR / f"targets_{month}.json"
    if not path.exists():
        raise FileNotFoundError(f"No mock target file for month {month}: {path}")
    return json.loads(path.read_text(encoding="utf-8"))


def _credential_summary(credential: Any) -> dict[str, Any]:
    if not isinstance(credential, dict):
        return {"present": False}
    secret = credential.get("secret")
    return {
        "present": True,
        "credential_ref": credential.get("credential_ref"),
        "credential_type": credential.get("credential_type"),
        "scope": credential.get("scope"),
        "secret_keys": sorted(secret.keys()) if isinstance(secret, dict) else [],
        "secret_returned": False,
    }


def _orders_summary(month: str) -> dict[str, Any]:
    rows = _read_csv(DATA_DIR / f"orders_{month}.csv")
    total = sum(float(row["amount"]) for row in rows)
    by_region: dict[str, float] = defaultdict(float)
    by_product: dict[str, float] = defaultdict(float)
    by_customer: dict[str, float] = defaultdict(float)
    for row in rows:
        amount = float(row["amount"])
        by_region[row["region"]] += amount
        by_product[row["product"]] += amount
        by_customer[row["customer"]] += amount
    return {
        "month": month,
        "source": "mock ERP orders",
        "row_count": len(rows),
        "revenue": round(total, 2),
        "top_region": max(by_region.items(), key=lambda x: x[1]),
        "top_product": max(by_product.items(), key=lambda x: x[1]),
        "top_customers": sorted(by_customer.items(), key=lambda x: x[1], reverse=True)[:5],
        "by_region": dict(sorted(by_region.items())),
        "by_product": dict(sorted(by_product.items())),
    }


def _finance_summary(month: str) -> dict[str, Any]:
    rows = _read_csv(DATA_DIR / f"invoices_{month}.csv")
    total = sum(float(row["amount"]) for row in rows)
    pending = sum(float(row["amount"]) for row in rows if row["status"] == "pending")
    by_category: dict[str, float] = defaultdict(float)
    by_supplier: dict[str, float] = defaultdict(float)
    for row in rows:
        amount = float(row["amount"])
        by_category[row["category"]] += amount
        by_supplier[row["supplier"]] += amount
    targets = _read_targets(month)
    watchlist = [
        {"supplier": row["supplier"], "invoice_id": row["invoice_id"], "amount": float(row["amount"]), "status": row["status"]}
        for row in rows
        if row["supplier"] in set(targets.get("watchlist_suppliers", []))
    ]
    return {
        "month": month,
        "source": "mock Finance invoices",
        "row_count": len(rows),
        "expense_total": round(total, 2),
        "pending_payables": round(pending, 2),
        "top_supplier": max(by_supplier.items(), key=lambda x: x[1]),
        "by_category": dict(sorted(by_category.items())),
        "watchlist_invoices": watchlist,
    }


def _write_csv(path: Path, rows: list[dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    keys: list[str] = []
    for row in rows:
        for key in row.keys():
            if key not in keys:
                keys.append(key)
    with path.open("w", encoding="utf-8", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=keys)
        writer.writeheader()
        writer.writerows(rows)


def _generate_html(path: Path, month: str, orders: dict[str, Any], finance: dict[str, Any], targets: dict[str, Any], insights: list[str]) -> None:
    revenue = orders["revenue"]
    expenses = finance["expense_total"]
    margin = revenue - expenses
    html = f"""<!doctype html>
<html><head><meta charset='utf-8'><title>LandGod Business Report Demo {month}</title>
<style>
body{{font-family:Inter,Arial,sans-serif;max-width:980px;margin:40px auto;color:#172033;line-height:1.5}}
.badge{{display:inline-block;background:#eef6ff;color:#075985;padding:4px 10px;border-radius:999px;font-size:12px}}
.grid{{display:grid;grid-template-columns:repeat(3,1fr);gap:14px;margin:22px 0}}
.card{{border:1px solid #dbe3ef;border-radius:14px;padding:16px;background:#fbfdff}}
.value{{font-size:28px;font-weight:700}}
li{{margin:8px 0}}
table{{border-collapse:collapse;width:100%;margin:16px 0}}td,th{{border:1px solid #dbe3ef;padding:8px;text-align:left}}th{{background:#f5f8fc}}
</style></head><body>
<span class='badge'>LandGod / MCPHub Enterprise Execution Harness Demo</span>
<h1>Monthly Business Review — {month}</h1>
<p>{targets.get('executive_theme','')}</p>
<div class='grid'>
<div class='card'><div>Revenue</div><div class='value'>${revenue:,.0f}</div><small>Target ${targets['revenue_target']:,.0f}</small></div>
<div class='card'><div>Expenses</div><div class='value'>${expenses:,.0f}</div><small>Budget ${targets['expense_budget']:,.0f}</small></div>
<div class='card'><div>Operating Margin</div><div class='value'>${margin:,.0f}</div><small>{margin/revenue:.1%} of revenue</small></div>
</div>
<h2>Executive Insights</h2><ul>{''.join(f'<li>{i}</li>' for i in insights)}</ul>
<h2>Revenue by Region</h2><table><tr><th>Region</th><th>Revenue</th></tr>{''.join(f'<tr><td>{k}</td><td>${v:,.0f}</td></tr>' for k,v in orders['by_region'].items())}</table>
<h2>Expense by Category</h2><table><tr><th>Category</th><th>Amount</th></tr>{''.join(f'<tr><td>{k}</td><td>${v:,.0f}</td></tr>' for k,v in finance['by_category'].items())}</table>
<h2>Audit Story</h2>
<p>Agent uses credential_ref only. Gateway issues a task-scoped grant. Worker executes the trusted business-report-demo tool locally. Gateway, Worker, and Credential audit records can be shown in WebUI.</p>
</body></html>"""
    path.write_text(html, encoding="utf-8")


def _generate_pptx(path: Path, month: str, orders: dict[str, Any], finance: dict[str, Any], targets: dict[str, Any], insights: list[str]) -> dict[str, Any]:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as exc:  # pragma: no cover - environment dependent
        return {"created": False, "reason": f"python-pptx unavailable: {exc}"}

    prs = Presentation()

    def add_title(slide, title, subtitle=None):
        slide.shapes.title.text = title
        if subtitle is not None:
            slide.placeholders[1].text = subtitle

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(slide, f"Monthly Business Review — {month}", "Generated by LandGod / MCPHub business-report-demo")

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Executive Scorecard"
    metrics = [
        ("Revenue", f"${orders['revenue']:,.0f}", f"Target ${targets['revenue_target']:,.0f}"),
        ("Expenses", f"${finance['expense_total']:,.0f}", f"Budget ${targets['expense_budget']:,.0f}"),
        ("Margin", f"${orders['revenue'] - finance['expense_total']:,.0f}", f"{(orders['revenue'] - finance['expense_total']) / orders['revenue']:.1%}"),
    ]
    for i, (label, value, note) in enumerate(metrics):
        box = slide.shapes.add_textbox(Inches(0.6 + i * 3.1), Inches(1.5), Inches(2.8), Inches(1.5))
        tf = box.text_frame
        tf.text = label
        p = tf.add_paragraph(); p.text = value; p.font.size = Pt(24); p.font.bold = True
        p = tf.add_paragraph(); p.text = note; p.font.size = Pt(12)

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Insights"
    slide.placeholders[1].text = "\n".join(f"• {x}" for x in insights)

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Enterprise Trust Story"
    slide.placeholders[1].text = "\n".join([
        "• Agent passes credential_ref, not raw secrets",
        "• Gateway issues a single-use task-scoped grant",
        "• Worker executes trusted tool locally",
        "• Gateway central audit + Worker local audit + Credential audit",
    ])

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    return {"created": True, "path": str(path), "slides": len(prs.slides)}


def _build_report(month: str, output_dir: str | None = None, credential: Any = None) -> dict[str, Any]:
    out_dir = Path(output_dir).expanduser().resolve() if output_dir else DEFAULT_OUTPUT_DIR.resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    orders = _orders_summary(month)
    finance = _finance_summary(month)
    targets = _read_targets(month)
    revenue = orders["revenue"]
    expenses = finance["expense_total"]
    margin = revenue - expenses
    revenue_delta = revenue - float(targets["revenue_target"])
    expense_delta = expenses - float(targets["expense_budget"])
    insights = [
        f"Revenue is ${revenue_delta:,.0f} {'above' if revenue_delta >= 0 else 'below'} target ({revenue / targets['revenue_target']:.1%} attainment).",
        f"Operating margin is ${margin:,.0f}, or {margin / revenue:.1%} of revenue.",
        f"Top revenue region is {orders['top_region'][0]} (${orders['top_region'][1]:,.0f}).",
        f"Top supplier spend is {finance['top_supplier'][0]} (${finance['top_supplier'][1]:,.0f}).",
        f"Expense budget variance is ${expense_delta:,.0f}; watchlist invoice count: {len(finance['watchlist_invoices'])}.",
    ]
    scorecard_rows = [
        {"metric": "revenue", "value": revenue, "target": targets["revenue_target"], "variance": revenue_delta},
        {"metric": "expenses", "value": expenses, "target": targets["expense_budget"], "variance": expense_delta},
        {"metric": "operating_margin", "value": margin, "target": round(float(targets["gross_margin_target"]) * revenue, 2), "variance": round(margin - float(targets["gross_margin_target"]) * revenue, 2)},
    ]
    scorecard_csv = out_dir / f"business_scorecard_{month}.csv"
    _write_csv(scorecard_csv, scorecard_rows)
    summary_json = out_dir / f"business_summary_{month}.json"
    summary = {
        "month": month,
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "credential": _credential_summary(credential),
        "orders": orders,
        "finance": finance,
        "targets": targets,
        "insights": insights,
        "artifacts": {},
    }
    html_path = out_dir / f"business_report_{month}.html"
    pptx_path = out_dir / f"business_report_{month}.pptx"
    _generate_html(html_path, month, orders, finance, targets, insights)
    pptx_status = _generate_pptx(pptx_path, month, orders, finance, targets, insights)
    audit_md = out_dir / f"audit_story_{month}.md"
    audit_md.write_text(
        "# LandGod Business Demo Audit Story\n\n"
        "- Agent only passes `credential_ref` to Gateway.\n"
        "- Gateway policy checks agent, worker group, tool, and scope.\n"
        "- Gateway issues a single-use task-scoped grant.\n"
        "- Worker exchanges the grant and injects `_landgod_credential` only into this trusted demo tool.\n"
        "- Gateway central audit, Worker local audit, and Credential audit prove the execution chain.\n",
        encoding="utf-8",
    )
    summary["artifacts"] = {
        "summary_json": str(summary_json),
        "scorecard_csv": str(scorecard_csv),
        "html_report": str(html_path),
        "pptx_report": str(pptx_path) if pptx_status.get("created") else None,
        "pptx_status": pptx_status,
        "audit_story": str(audit_md),
    }
    summary_json.write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    return summary


def _tool_schema(include_output_dir: bool = False, include_credential: bool = False) -> dict[str, Any]:
    properties: dict[str, Any] = {
        "month": {"type": "string", "default": "2026-06", "description": "Reporting month, e.g. 2026-06"},
    }
    if include_output_dir:
        properties["output_dir"] = {"type": "string", "description": "Directory where report artifacts should be written"}
    if include_credential:
        properties["_landgod_credential"] = {"type": "object", "description": "Injected by LandGod Credential Broker; callers should pass credential_ref, not this field"}
    return {"type": "object", "properties": properties, "additionalProperties": True}


@server.list_tools()
async def list_tools() -> list[Tool]:
    return [
        Tool(
            name="load_erp_orders",
            description="Read mock ERP order data and return a revenue summary. Credential-aware; never returns secret values.",
            inputSchema=_tool_schema(include_credential=True),
        ),
        Tool(
            name="load_finance_invoices",
            description="Read mock Finance invoice data and return spend/watchlist summary. Credential-aware; never returns secret values.",
            inputSchema=_tool_schema(include_credential=True),
        ),
        Tool(
            name="generate_monthly_report",
            description="Generate CSV, HTML, PPTX, and audit-story artifacts from mock ERP/Finance data.",
            inputSchema=_tool_schema(include_output_dir=True),
        ),
        Tool(
            name="run_monthly_close_demo",
            description="End-to-end business demo: read ERP + Finance mock data, verify injected credential metadata, and generate executive artifacts.",
            inputSchema=_tool_schema(include_output_dir=True, include_credential=True),
        ),
    ]


@server.call_tool(validate_input=False)
async def call_tool(name: str, arguments: dict[str, Any]) -> dict[str, Any]:
    args = arguments or {}
    month = str(args.get("month") or "2026-06")
    credential = args.get("_landgod_credential")
    if name == "load_erp_orders":
        result = _orders_summary(month)
        result["credential"] = _credential_summary(credential)
        return result
    if name == "load_finance_invoices":
        result = _finance_summary(month)
        result["credential"] = _credential_summary(credential)
        return result
    if name == "generate_monthly_report":
        return _build_report(month, args.get("output_dir"), None)
    if name == "run_monthly_close_demo":
        return _build_report(month, args.get("output_dir"), credential)
    raise ValueError(f"Unknown tool: {name}")


async def _run_stdio() -> None:
    async with stdio_server() as (read_stream, write_stream):
        await server.run(read_stream, write_stream, server.create_initialization_options())


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--smoke", action="store_true", help="Run the demo once without MCP stdio.")
    parser.add_argument("--month", default="2026-06")
    parser.add_argument("--output-dir", default=str(DEFAULT_OUTPUT_DIR))
    args = parser.parse_args()
    if args.smoke:
        print(json.dumps(_build_report(args.month, args.output_dir, {"credential_ref": "smoke", "credential_type": "api_token", "scope": "report", "secret": {"token": "not-returned"}}), ensure_ascii=False, indent=2))
        return
    anyio.run(_run_stdio)


if __name__ == "__main__":
    main()
